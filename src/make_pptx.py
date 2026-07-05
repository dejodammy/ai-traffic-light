"""Build a modern, visual Final Year Project defense deck."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "results", "figures")

NAVY = RGBColor(0x0E, 0x1F, 0x56)
NAVY2 = RGBColor(0x1B, 0x32, 0x7A)
GREEN = RGBColor(0x1F, 0x9D, 0x55)
AMBER = RGBColor(0xF2, 0xA8, 0x2B)
RED = RGBColor(0xD6, 0x3B, 0x3B)
DARK = RGBColor(0x22, 0x2A, 0x33)
GREY = RGBColor(0x6B, 0x76, 0x86)
BG = RGBColor(0xF1, 0xF4, 0xF9)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTE = RGBColor(0xC7, 0xD0, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
ACCENTS = [GREEN, AMBER, NAVY2, RED, GREEN, AMBER]


def slide():
    return prs.slides.add_slide(BLANK)


def shape(s, kind, x, y, w, h, fill=None, line=None, line_w=1.25, round_=0.09, shadow=False):
    sp = s.shapes.add_shape(kind, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = round_
        except Exception: pass
    return sp


def txt(s, x, y, w, h, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri", spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = line
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb


def bg(s, color=BG):
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=color)


def tlight(s, x, y, w, dark=NAVY):
    """traffic-light motif: dark rounded body + 3 lamps."""
    h = w * 2.6
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=dark, round_=0.3)
    d = w * 0.56; gap = (h - 3 * d) / 4
    for i, c in enumerate([RED, AMBER, GREEN]):
        shape(s, MSO_SHAPE.OVAL, x + (w - d) / 2, y + gap * (i + 1) + d * i, d, d, fill=c)


def header(s, title, idx):
    bg(s)
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.18), fill=NAVY)
    shape(s, MSO_SHAPE.RECTANGLE, 0, Inches(1.18), SW, Inches(0.09), fill=GREEN)
    # mini traffic-light dots
    for i, c in enumerate([RED, AMBER, GREEN]):
        shape(s, MSO_SHAPE.OVAL, Inches(0.55), Inches(0.30 + i * 0.205), Inches(0.16), Inches(0.16), fill=c)
    txt(s, Inches(1.0), 0, Inches(11.6), Inches(1.18), title, 28, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # footer
    shape(s, MSO_SHAPE.OVAL, Inches(12.55), Inches(6.92), Inches(0.5), Inches(0.5), fill=NAVY)
    txt(s, Inches(12.55), Inches(6.92), Inches(0.5), Inches(0.5), str(idx), 14, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.55), Inches(6.98), Inches(9), Inches(0.4),
        "Intelligent Traffic Light System · Deep Reinforcement Learning", 10.5, GREY)


def feature(s, x, y, w, idx, title, desc=None, color=None):
    color = color or ACCENTS[idx % len(ACCENTS)]
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.95 if desc else 0.74),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.14)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), Inches(0.95 if desc else 0.74),
          fill=color, round_=0.5)
    shape(s, MSO_SHAPE.OVAL, x + Inches(0.28), y + Inches(0.19), Inches(0.36), Inches(0.36), fill=color)
    txt(s, x + Inches(0.28), y + Inches(0.19), Inches(0.36), Inches(0.36), str(idx + 1), 13, WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.85), y + (Inches(0.10) if desc else 0), w - Inches(1.0),
        Inches(0.45 if desc else 0.74), title, 16.5, DARK, bold=True,
        anchor=MSO_ANCHOR.MIDDLE if not desc else MSO_ANCHOR.TOP)
    if desc:
        txt(s, x + Inches(0.85), y + Inches(0.5), w - Inches(1.0), Inches(0.4), desc, 12.5, GREY)


def stat(s, x, y, w, h, value, label, color):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2),
          line_w=1.0, round_=0.12)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.14), fill=color, round_=0.5)
    txt(s, x, y + Inches(0.30), w, Inches(0.9), value, 40, color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y + h - Inches(0.62), w, Inches(0.5), label, 13, GREY, align=PP_ALIGN.CENTER, bold=True)


def image(s, path, x, y, w=None, h=None):
    p = os.path.join(FIG, path)
    if os.path.exists(p):
        s.shapes.add_picture(p, x, y, width=w, height=h)


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


# ============================ 1. TITLE ============================
s = slide()
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=NAVY)
shape(s, MSO_SHAPE.RECTANGLE, Inches(9.1), 0, Inches(4.23), SH, fill=NAVY2)       # right panel
shape(s, MSO_SHAPE.RECTANGLE, Inches(9.1), 0, Inches(0.10), SH, fill=GREEN)
tlight(s, Inches(10.7), Inches(2.0), Inches(1.05))
txt(s, Inches(0.7), Inches(1.2), Inches(8.1), Inches(0.5), "FINAL YEAR PROJECT  ·  DEFENSE",
    15, GREEN, bold=True)
txt(s, Inches(0.7), Inches(1.8), Inches(8.2), Inches(2.6),
    "An Intelligent Traffic Light System Using Deep Reinforcement Learning",
    36, WHITE, bold=True, spacing=1.05)
txt(s, Inches(0.7), Inches(4.35), Inches(8.2), Inches(0.6), "A Case Study of Lagos Traffic",
    20, MUTE, italic=True)
shape(s, MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(5.05), Inches(2.2), Inches(0.05), fill=GREEN)
txt(s, Inches(0.7), Inches(5.25), Inches(8.2), Inches(0.5),
    "Oluwadamilola Samson Oladejo   ·   21120612543", 18, WHITE, bold=True)
txt(s, Inches(0.7), Inches(5.8), Inches(8.2), Inches(1.4),
    "Computer & Information Sciences · Pan-Atlantic University\n"
    "Supervisor: [Supervisor's Name]            July 2026", 14, MUTE, spacing=1.2)
notes(s, "Good day, panel. My name is Oluwadamilola Oladejo. My project is an intelligent traffic "
         "light system that uses deep reinforcement learning to control signals adaptively, with a "
         "busy junction in Ikeja, Lagos, as the case study. I will cover the problem, my approach, my "
         "results, and the honest conclusions.")

# ============================ 2. INTRODUCTION ============================
s = slide(); header(s, "Introduction — Background", 2)
stat(s, Inches(0.6), Inches(1.55), Inches(3.9), Inches(1.7), "#1", "Lagos: among the world's\nmost congested cities", RED)
stat(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(1.7), "~4 hrs", "lost daily by commuters\nin traffic", AMBER)
stat(s, Inches(8.8), Inches(1.55), Inches(3.9), Inches(1.7), "₦4tn", "estimated annual loss\nin productivity", NAVY2)
feature(s, Inches(0.6), Inches(3.55), Inches(12.1), 0, "Most Lagos junctions use FIXED-TIME signals — pre-set timers that ignore live traffic.")
feature(s, Inches(0.6), Inches(4.45), Inches(12.1), 1, "Congestion drives higher fuel use, emissions and air pollution.")
feature(s, Inches(0.6), Inches(5.35), Inches(12.1), 2, "AI and reinforcement learning now make adaptive, self-learning signal control feasible.")
notes(s, "Lagos is consistently ranked among the world's worst cities for traffic. Commuters lose "
         "hours daily and the economy loses trillions of naira, with serious pollution effects. A key "
         "reason is that most signals are fixed-time — they cannot react to actual traffic. AI and "
         "reinforcement learning now offer a smarter, adaptive alternative, and that is the "
         "opportunity this project explores.")

# ============================ 3. PROBLEM ============================
s = slide(); header(s, "Problem Statement", 3)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.5),
      fill=RGBColor(0xFC, 0xEC, 0xEC), line=RED, line_w=1.5, round_=0.06)
txt(s, Inches(1.0), Inches(1.55), Inches(11.3), Inches(1.5),
    "A fixed-time signal is BLIND: it gives the same green whether a road is empty or jammed.",
    22, RED, bold=True, anchor=MSO_ANCHOR.MIDDLE)
feature(s, Inches(0.6), Inches(3.35), Inches(12.1), 0, "Unnecessary delay off-peak; worse congestion at peak hours.")
feature(s, Inches(0.6), Inches(4.25), Inches(12.1), 1, "No response to changing demand and no priority for emergency vehicles.")
feature(s, Inches(0.6), Inches(5.15), Inches(12.1), 2, "Need: an adaptive controller driven by real-time conditions — on affordable hardware.")
notes(s, "The core problem is that a fixed timer is blind. It gives the same green to an empty road as "
         "to a jammed one. This wastes time off-peak, worsens jams at peak, and cannot prioritise an "
         "ambulance. So we need an adaptive controller that decides from live conditions, and ideally "
         "one that runs on low-cost hardware.")

# ============================ 4. AIM & OBJECTIVES ============================
s = slide(); header(s, "Aim and Objectives", 4)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0),
      fill=NAVY, round_=0.08)
txt(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0),
    "AIM:  Model and evaluate an intelligent traffic light system that adjusts signal timing from "
    "real-time conditions using reinforcement learning.", 16, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
objs = ["Review fixed-time and AI-based traffic light systems.",
        "Design a traffic state representation and reward function.",
        "Implement a Deep Q-Network (DQN) controller in simulation.",
        "Build a realistic SUMO model of a real intersection.",
        "Integrate camera-based vehicle detection (YOLO)."]
for i, o in enumerate(objs):
    feature(s, Inches(0.6), Inches(2.75 + i * 0.78), Inches(12.1), i, o)
notes(s, "The aim was to model and evaluate an adaptive, RL-based signal system. I set five "
         "objectives: review the field; design the state and reward; implement a DQN controller; "
         "build a realistic SUMO model of a real junction; and integrate YOLO vehicle detection so "
         "the same controller can run on a camera feed.")

# ============================ 5. RESEARCH QUESTIONS ============================
s = slide(); header(s, "Research Questions", 5)
qs = [("Can a DQN reduce delay and queues versus fixed-time control at a real junction?", GREEN),
      ("How does it compare against an OPTIMISED timer (Webster) and an ACTUATED controller?", AMBER),
      ("Does the learned policy GENERALISE to different and unseen demand patterns?", NAVY2),
      ("Under which conditions is reinforcement-learning control most beneficial?", RED)]
for i, (q, c) in enumerate(qs):
    y = Inches(1.6 + i * 1.28)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(1.08),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.1)
    shape(s, MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.24), Inches(0.6), Inches(0.6), fill=c)
    txt(s, Inches(0.85), y + Inches(0.24), Inches(0.6), Inches(0.6), f"Q{i+1}", 16, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.7), y, Inches(10.8), Inches(1.08), q, 17, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Four questions guided the work. Does RL beat a fixed timer at a real junction? How does it "
         "compare to an optimised Webster timer and an actuated controller — not just a naive one? "
         "Does it generalise rather than memorise? And where exactly is it worth using? I answered "
         "the last two honestly.")

# ============================ 6. LITERATURE ============================
s = slide(); header(s, "Literature Review", 6)
cols = [("Fixed-time & Webster", "Simple, widely used — but blind to real-time demand.", GREY),
        ("Actuated control", "Sensor-based; serve the longest queue. Strong but rule-based.", AMBER),
        ("Reinforcement learning", "DQN, PressLight, PDLight — state of the art, adaptive.", GREEN)]
for i, (t, d, c) in enumerate(cols):
    x = Inches(0.6 + i * 4.1)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.55), Inches(3.85), Inches(2.3),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.08)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.55), Inches(3.85), Inches(0.55), fill=c, round_=0.12)
    txt(s, x, Inches(1.55), Inches(3.85), Inches(0.55), t, 15, WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.25), Inches(2.3), Inches(3.35), Inches(1.4), d, 13.5, DARK)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.15), Inches(12.1), Inches(2.4),
      fill=RGBColor(0xEC, 0xF4, 0xEF), line=GREEN, line_w=1.5, round_=0.05)
txt(s, Inches(0.95), Inches(4.32), Inches(11.4), Inches(0.5), "RESEARCH GAP", 15, GREEN, bold=True)
txt(s, Inches(0.95), Inches(4.8), Inches(11.4), Inches(1.6),
    "Most RL studies use abstract/synthetic networks and weak baselines. Few validate on a real, "
    "geometry-faithful Lagos junction against STRONG baselines (Webster + actuated), using an "
    "accumulated-delay metric that reflects real driver experience.", 16, DARK, spacing=1.1)
notes(s, "Classical control spans fixed-time, Webster timing, and actuated control. On the AI side, "
         "DQN and pressure-based methods like PressLight and PDLight are state of the art. But most RL "
         "papers test on synthetic grids against weak baselines. My gap is a rigorous, honest "
         "evaluation on a real, validated Lagos junction, against strong baselines, with a realistic "
         "delay metric.")

# ============================ 7. METHODOLOGY ============================
s = slide(); header(s, "Methodology & System Architecture", 7)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.45), Inches(8.35), Inches(3.4),
      fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.04)
image(s, "architecture.png", Inches(0.6), Inches(1.6), w=Inches(8.05))
items = [("SUMO + TraCI", "microscopic simulator", GREEN),
         ("Real Ikeja junction", "OpenStreetMap, validated vs satellite", NAVY2),
         ("YOLO vision", "camera-based vehicle detection", AMBER),
         ("Double DQN", "+ prioritised experience replay", GREEN),
         ("3 baselines", "fixed-time, Webster, actuated", RED)]
for i, (t, d, c) in enumerate(items):
    y = Inches(1.5 + i * 1.02)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), y, Inches(3.85), Inches(0.9),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.12)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), y, Inches(0.12), Inches(0.9), fill=c, round_=0.5)
    txt(s, Inches(9.2), y + Inches(0.08), Inches(3.5), Inches(0.4), t, 14.5, DARK, bold=True)
    txt(s, Inches(9.2), y + Inches(0.46), Inches(3.5), Inches(0.4), d, 11.5, GREY)
notes(s, "My method combines a microscopic simulator (SUMO via TraCI); a real Ikeja junction exported "
         "from OpenStreetMap and corrected against satellite imagery; YOLO for vehicle detection; and "
         "a Double DQN controller. As the diagram shows, the camera or simulator produces a state, the "
         "DQN chooses a phase, and that drives the signal — a closed feedback loop. I compare against "
         "three baselines, including the strong ones.")

# ============================ 8. SYSTEM DESIGN ============================
s = slide(); header(s, "System Design & Implementation", 8)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(1.5), Inches(5.1), Inches(3.0),
      fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.05)
image(s, "rl_loop.png", Inches(7.85), Inches(1.65), w=Inches(4.8))
specs = [("STATE", "14 features — queue, accumulated wait, density / approach + phase", GREEN),
         ("ACTION", "choose the green phase (2-phase or 4-phase signal)", AMBER),
         ("REWARD", "−(queue + weighted accumulated wait) − switch penalty", NAVY2),
         ("LEARNING", "ε-greedy · target network · γ = 0.99", RED)]
for i, (t, d, c) in enumerate(specs):
    y = Inches(1.55 + i * 1.0)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y, Inches(6.9), Inches(0.88),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.1)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), y, Inches(1.45), Inches(0.88), fill=c, round_=0.1)
    txt(s, Inches(0.55), y, Inches(1.45), Inches(0.88), t, 13, WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.15), y, Inches(5.1), Inches(0.88), d, 13, DARK, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(7.7), Inches(4.7), Inches(5.1), Inches(0.5), "Stack: Python · PyTorch · YOLO · SUMO · Raspberry Pi/GPIO",
    13, GREY, align=PP_ALIGN.CENTER, italic=True)
notes(s, "The agent observes 14 numbers — queue, accumulated waiting time and density for each "
         "approach, plus the current phase. Its action is which phase to make green. The reward "
         "penalises queues and accumulated delay, with a small switch penalty to avoid unsafe flicker. "
         "It is built in Python with PyTorch, YOLO and SUMO, and designed to deploy on a Raspberry Pi "
         "driving real lights via GPIO.")

# ============================ 9. RESULTS 1 ============================
s = slide(); header(s, "Results — Two-Phase Control", 9)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.4), Inches(8.5), Inches(4.9),
      fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.04)
image(s, "baseline_comparison.png", Inches(0.6), Inches(2.5), w=Inches(8.2))
txt(s, Inches(0.6), Inches(1.55), Inches(8.2), Inches(0.9),
    "On a conventional two-phase signal, the DQN beat EVERY baseline — lower queues, lower delay, higher throughput.",
    14, DARK, bold=True)
stat(s, Inches(9.15), Inches(1.5), Inches(3.6), Inches(1.5), "−69%", "delay vs Naive fixed-time", GREEN)
stat(s, Inches(9.15), Inches(3.15), Inches(3.6), Inches(1.5), "−60%", "delay vs Webster-optimal", NAVY2)
stat(s, Inches(9.15), Inches(4.8), Inches(3.6), Inches(1.5), "−76%", "delay vs Actuated", AMBER)
notes(s, "Here are the headline two-phase results. The chart shows queue, delay and throughput per "
         "controller at rush hour. The DQN has the lowest queue and delay and the highest throughput. "
         "It cuts accumulated delay by 69 percent against the naive timer, 60 percent against the "
         "optimised Webster timer, and 76 percent against the actuated controller — and it wins on all "
         "demand patterns, not just this one.")

# ============================ 10. RESULTS 2 ============================
s = slide(); header(s, "Results — Effect of Signal Architecture", 10)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.45), Inches(8.3), Inches(4.6),
      fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.04)
image(s, "dqn_vs_actuated.png", Inches(0.6), Inches(1.7), w=Inches(8.0))
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.55), Inches(3.85), Inches(2.15),
      fill=RGBColor(0xEC, 0xF4, 0xEF), line=GREEN, line_w=1.3, round_=0.07)
txt(s, Inches(9.2), Inches(1.7), Inches(3.4), Inches(0.5), "TWO-PHASE", 14, GREEN, bold=True)
txt(s, Inches(9.2), Inches(2.15), Inches(3.4), Inches(1.5),
    "Beats fixed-time, Webster AND actuated — across all demand.", 14, DARK)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(3.9), Inches(3.85), Inches(2.4),
      fill=RGBColor(0xFD, 0xF3, 0xE3), line=AMBER, line_w=1.3, round_=0.07)
txt(s, Inches(9.2), Inches(4.05), Inches(3.4), Inches(0.5), "FOUR-PHASE", 14, AMBER, bold=True)
txt(s, Inches(9.2), Inches(4.5), Inches(3.4), Inches(1.7),
    "Beats fixed-time & Webster — but a near-optimal actuated rule stays ahead. Even PressLight & "
    "PDLight didn't beat it here. (Honest, literature-consistent.)", 13, DARK, spacing=1.05)
notes(s, "This slide is about honest nuance. On two-phase signals — green bars — the DQN wins "
         "everywhere. On four-phase protected signals — amber bars — it still beats fixed-time and "
         "Webster, but actuated control is better. I even implemented PressLight and PDLight, and they "
         "did not beat actuated at a single junction either, which matches the literature. Reporting "
         "this boundary makes the result trustworthy.")

# ============================ 11. DISCUSSION ============================
s = slide(); header(s, "Discussion", 11)
feature(s, Inches(0.6), Inches(1.55), Inches(12.1), 0,
        "The DQN's edge is REAL-TIME RESPONSIVENESS — it never wastes green on an empty approach.")
feature(s, Inches(0.6), Inches(2.45), Inches(12.1), 1,
        "Gains are largest under imbalanced / variable demand; smallest under balanced, steady demand.")
feature(s, Inches(0.6), Inches(3.35), Inches(12.1), 2,
        "For protected 4-phase junctions, a simple actuated rule is already near-optimal.")
feature(s, Inches(0.6), Inches(4.25), Inches(12.1), 3,
        "Findings align with traffic-engineering theory and meet all five objectives.")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.1),
      fill=NAVY, round_=0.08)
txt(s, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.1),
    "Deployment guidance: use RL on conventional two-phase junctions under heavy, imbalanced Lagos peak demand.",
    17, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Why does it work? The agent reacts to the live state and never wastes green on an empty "
         "road, so its gains are biggest where fixed timers struggle — imbalanced, variable demand. "
         "Where demand is balanced or the signal is a saturated four-phase one, the actuated rule is "
         "already near-optimal. This is consistent with theory and gives a clear deployment message: "
         "use RL on two-phase junctions under heavy, imbalanced peak demand — exactly the Lagos case.")

# ============================ 12. CHALLENGES ============================
s = slide(); header(s, "Challenges & Limitations", 12)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.55),
      fill=RGBColor(0xFD, 0xF3, 0xE3), line=AMBER, line_w=1.5, round_=0.05)
txt(s, Inches(0.95), Inches(1.62), Inches(11.4), Inches(0.5), "KEY CHALLENGE — POLICY COLLAPSE", 15, AMBER, bold=True)
txt(s, Inches(0.95), Inches(2.08), Inches(11.4), Inches(0.95),
    "The agent first learned to serve only one approach. Diagnosed as a state–reward mismatch "
    "(partial observability) and fixed by exposing accumulated wait in the state.", 15, DARK, spacing=1.05)
feature(s, Inches(0.6), Inches(3.3), Inches(12.1), 0, "Simulation-based; demand estimated, not measured at the junction.")
feature(s, Inches(0.6), Inches(4.2), Inches(12.1), 1, "Single junction — corridor (multi-junction) coordination not addressed.")
feature(s, Inches(0.6), Inches(5.1), Inches(12.1), 2, "Actuated superior on 4-phase; YOLO field accuracy not yet quantified.")
notes(s, "The main challenge was a policy collapse: the agent served only one approach and starved the "
         "other. I traced it to a subtle flaw — the reward penalised accumulated delay, but the state "
         "only showed instantaneous delay, so the agent was punished for something it could not see. "
         "Exposing accumulated delay in the state fixed it. Honest limitations: simulation-based with "
         "estimated demand; a single junction; actuated wins on four-phase; and YOLO accuracy is not "
         "yet field-measured.")

# ============================ 13. CONCLUSION ============================
s = slide(); header(s, "Conclusion", 13)
stat(s, Inches(0.6), Inches(1.55), Inches(3.9), Inches(1.9), "16–76%", "delay cut on TWO-PHASE\nvs all baselines", GREEN)
stat(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(1.9), "4 / 4", "demand patterns won\non two-phase", NAVY2)
stat(s, Inches(8.8), Inches(1.55), Inches(3.9), Inches(1.9), "Real", "Ikeja junction from\nOpenStreetMap", AMBER)
feature(s, Inches(0.6), Inches(3.75), Inches(12.1), 0, "Two-phase: DQN beat fixed-time, Webster AND actuated on every demand pattern.")
feature(s, Inches(0.6), Inches(4.65), Inches(12.1), 1, "Four-phase: beat fixed-time and Webster; a gap-out actuated controller remained best.")
feature(s, Inches(0.6), Inches(5.55), Inches(12.1), 2, "Most valuable for two-phase junctions under heavy, imbalanced Lagos peak traffic.")
notes(s, "To conclude: on two-phase control my agent cut accumulated delay by 16 to 76 percent against "
         "every baseline, on every demand pattern. On four-phase it beat the fixed timers while "
         "actuated remained best. Its advantage is responsiveness, and the honest takeaway is that RL "
         "is most valuable for two-phase junctions under the heavy, imbalanced demand of Lagos rush "
         "hour.")

# ============================ 14. RECOMMENDATIONS ============================
s = slide(); header(s, "Recommendations & Future Work", 14)
recs = ["Collect field data — run the project's YOLO on real junction footage.",
        "Extend to corridor (multi-agent) control along Obafemi Awolowo Way.",
        "Add a throughput-aware / max-pressure term to the reward.",
        "Robust perception — fine-tune for okada and keke under Lagos conditions.",
        "Hardware pilot — supervised Raspberry Pi + GPIO with a safety fallback.",
        "Implement and evaluate emergency-vehicle preemption."]
for i, r in enumerate(recs):
    col = i % 2; row = i // 2
    feature(s, Inches(0.6 + col * 6.15), Inches(1.7 + row * 1.55), Inches(5.95), i, r)
notes(s, "For future work: use my own YOLO on real footage to get measured demand; extend to a "
         "corridor of signals where pressure-based RL should finally win; add a throughput term to the "
         "reward; harden the vision for motorcycles and tricycles; run a supervised hardware pilot; "
         "and complete emergency-vehicle priority, which directly addresses the original problem.")

# ============================ 15. REFERENCES ============================
s = slide(); header(s, "References", 15)
refs = ["Krajzewicz, D., et al. (2012). Recent development and applications of SUMO. Int. J. Adv. Systems & Measurements, 5(3–4), 128–138.",
        "Mnih, V., et al. (2015). Human-level control through deep reinforcement learning. Nature, 518, 529–533.",
        "Sutton, R. S., & Barto, A. G. (2018). Reinforcement Learning: An Introduction (2nd ed.). MIT Press.",
        "Webster, F. V. (1958). Traffic Signal Settings. Road Research Technical Paper No. 39. HMSO.",
        "Wei, H., et al. (2019). PressLight: Learning max pressure control to coordinate traffic signals. KDD 2019.",
        "Zhang, C., et al. (2020). PDLight: DRL traffic light control with pressure and dynamic duration. arXiv:2009.13711."]
for i, r in enumerate(refs):
    y = Inches(1.6 + i * 0.86)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(0.74),
          fill=CARD, line=RGBColor(0xE2, 0xE7, 0xF2), line_w=1.0, round_=0.1)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(0.12), Inches(0.74),
          fill=ACCENTS[i % len(ACCENTS)], round_=0.5)
    txt(s, Inches(0.95), y, Inches(11.6), Inches(0.74), r, 13, DARK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "These are the key references in APA format: the SUMO simulator, the DQN paper by Mnih and "
         "colleagues, Sutton and Barto's textbook, Webster's signal-timing method, and the two "
         "state-of-the-art methods I implemented — PressLight and PDLight.")

# ============================ 16. THANK YOU ============================
s = slide()
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=NAVY)
shape(s, MSO_SHAPE.RECTANGLE, Inches(9.1), 0, Inches(4.23), SH, fill=NAVY2)
shape(s, MSO_SHAPE.RECTANGLE, Inches(9.1), 0, Inches(0.10), SH, fill=GREEN)
tlight(s, Inches(10.75), Inches(2.6), Inches(1.0))
txt(s, Inches(0.8), Inches(2.5), Inches(8.0), Inches(1.4), "Thank You", 56, WHITE, bold=True)
shape(s, MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(3.95), Inches(2.4), Inches(0.06), fill=GREEN)
txt(s, Inches(0.8), Inches(4.2), Inches(8.0), Inches(0.7), "Questions & Discussion", 24, MUTE, italic=True)
txt(s, Inches(0.8), Inches(5.3), Inches(8.0), Inches(0.6),
    "Oluwadamilola Samson Oladejo", 20, WHITE, bold=True)
txt(s, Inches(0.8), Inches(5.85), Inches(8.0), Inches(0.5), "dejodammy7@gmail.com", 16, MUTE)
notes(s, "Thank you for your attention. I am happy to take questions, and I can walk through the live "
         "simulation, the training curves, or the diagnostic finding in more detail.")

out = os.path.join(ROOT, "Defense_Presentation_v2.pptx")
prs.save(out)
print("saved", out, "| slides:", len(prs.slides._sldIdLst))
